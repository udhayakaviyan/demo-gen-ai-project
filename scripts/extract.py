import pytesseract
from PIL import Image
import cv2
import os
from pptx import Presentation
import subprocess
import easyocr

ocr_reader = easyocr.Reader(['en'], gpu=False)

def extract_text_from_image(image_path):
            result = ocr_reader.readtext(image_path, detail=0)
            text = "\n".join(result).strip()
            return text

def extract_text_from_pptx(pptx_path):
    prs = Presentation(pptx_path)
    content = []

    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                content.append(shape.text)
            elif shape.shape_type == 19:  # Table
                for row in shape.table.rows:
                    row_text = [cell.text.strip() for cell in row.cells]
                    content.append("\t".join(row_text))
            # Extract images with OCR
            elif shape.shape_type == 13 and hasattr(shape, "image"):
                try:
                    image = shape.image
                    ext = image.ext
                    image_data = image.blob
                    image_filename = f"{shape.image.sha1}.{ext}"
                    image_path = os.path.join("output_images", image_filename)
                    with open(image_path, "wb") as f:
                        f.write(image_data)
                    result = ocr_reader.readtext(image_path, detail=0)
                    text = "\n".join(result).strip()
                 # text = extract_text_from_image(cleaned_path)
                    # print(text)
                    content.append(text)
                except Exception as e:
                    content.append(f"[Image OCR failed: {e}]")
    return "\n".join(content)

def extract_from_folder(image_folder, pptx_folder):
    docs = []
    if any(f.lower().endswith(".ppt") and os.path.isfile(os.path.join(pptx_folder, f)) for f in os.listdir(pptx_folder)):
        convert_ppt_to_pptx(pptx_folder)
    for filename in os.listdir(image_folder):
        if filename.lower().endswith((".png", ".jpg", ".jpeg")):
            path = os.path.join(image_folder, filename)
            text = extract_text_from_image(path)
            # print(text)
            #  print("text",text)
            docs.append({"filename": filename, "content": text})
    for filename in os.listdir(pptx_folder):
        if filename.lower().endswith(".pptx"):
            path = os.path.join(pptx_folder, filename)
            text = extract_text_from_pptx(path)
            # print(text)
            docs.append({"filename": filename, "content": text})      
    return docs

def convert_ppt_to_pptx(ppt_folder):
    for file in os.listdir(ppt_folder):
        if file.lower().endswith(".ppt"):
            input_path =os.path.abspath( os.path.join(ppt_folder, file))
            #"data/pptx\agile vs waterfall.ppt"
            assert os.path.exists(input_path)
            try:
                result =subprocess.run([r"C:\Program Files\LibreOffice\program\soffice.exe", "--headless", "--convert-to", "pptx","--outdir",ppt_folder, input_path],
                check=True,capture_output=True,text=True)
                
                print(f"✅ Converted: {file}")
            except subprocess.CalledProcessError as e:
                print(f"❌ Failed: {file} - {e}")